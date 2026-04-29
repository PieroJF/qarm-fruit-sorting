"""
Generate Final Presentation — Robotic Pick & Place for Fruit Sorting
Applied Robotics (04 39984) — University of Birmingham

Run:  C:/Python313/python.exe slides/generate_final_pptx.py
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
import os

# ── Colours (A3 palette) ──
DARK_BG      = RGBColor(0x1A, 0x1A, 0x2E)
WHITE        = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY   = RGBColor(0xCC, 0xCC, 0xCC)
ACCENT_BLUE  = RGBColor(0x00, 0x96, 0xD6)
ACCENT_ORANGE= RGBColor(0xFF, 0x8C, 0x00)
ACCENT_GREEN = RGBColor(0x4C, 0xAF, 0x50)
ACCENT_RED   = RGBColor(0xFF, 0x6B, 0x6B)
SECTION_BG   = RGBColor(0x22, 0x22, 0x3A)

TOTAL_SLIDES = 8
_ROOT = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)


# ── Helpers (matching A3 signatures) ──

def add_bg(slide, color=DARK_BG):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_textbox(slide, left, top, width, height, text, font_size=14,
                bold=False, color=WHITE, alignment=PP_ALIGN.LEFT,
                font_name='Calibri'):
    txBox = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = alignment
    return tf


def add_paragraph(tf, text, font_size=14, bold=False, color=WHITE,
                  font_name='Calibri', alignment=PP_ALIGN.LEFT,
                  space_before=0, space_after=0):
    p = tf.add_paragraph()
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = alignment
    if space_before:
        p.space_before = Pt(space_before)
    if space_after:
        p.space_after = Pt(space_after)
    return p


def add_rect(slide, left, top, width, height,
             fill_color=SECTION_BG, border_color=None):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(left), Inches(top), Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    return shape


def add_placeholder_box(slide, left, top, width, height,
                        label="[INSERT SCREENSHOT]"):
    shape = add_rect(slide, left, top, width, height,
                     RGBColor(0x33, 0x33, 0x55), ACCENT_BLUE)
    shape.text_frame.word_wrap = True
    p = shape.text_frame.paragraphs[0]
    p.text = label
    p.font.size = Pt(11)
    p.font.color.rgb = ACCENT_BLUE
    p.font.italic = True
    p.alignment = PP_ALIGN.CENTER
    shape.text_frame.paragraphs[0].space_before = Pt(
        max(0, int((height * 72 / 2.54) / 3)))
    return shape


def add_image_or_placeholder(slide, left, top, width, height,
                             rel_path, placeholder_label):
    abs_path = os.path.join(_ROOT, rel_path)
    if os.path.exists(abs_path):
        slide.shapes.add_picture(abs_path,
                                 Inches(left), Inches(top),
                                 Inches(width), Inches(height))
    else:
        add_placeholder_box(slide, left, top, width, height,
                            placeholder_label)


def header_bar(slide, title, subtitle=""):
    add_rect(slide, 0, 0, 13.333, 1.0, ACCENT_BLUE)
    add_textbox(slide, 0.4, 0.1, 9, 0.5, title,
                font_size=26, bold=True, color=WHITE)
    if subtitle:
        add_textbox(slide, 0.4, 0.55, 9, 0.35, subtitle,
                    font_size=13, color=RGBColor(0xDD, 0xEE, 0xFF))
    add_textbox(slide, 9.5, 0.1, 3.5, 0.5,
                "Final Project | Applied Robotics",
                font_size=10, color=RGBColor(0xDD, 0xEE, 0xFF),
                alignment=PP_ALIGN.RIGHT)


def footer_bar(slide, slide_num):
    add_rect(slide, 0, 7.1, 13.333, 0.4, RGBColor(0x11, 0x11, 0x22))
    add_textbox(slide, 0.4, 7.12, 8, 0.35,
                "Piero Flores | Zihen Huang | Ran Zhang | Yichang Chao",
                font_size=9, color=LIGHT_GRAY)
    add_textbox(slide, 11.5, 7.12, 1.5, 0.35,
                f"Slide {slide_num}/{TOTAL_SLIDES}",
                font_size=9, color=LIGHT_GRAY, alignment=PP_ALIGN.RIGHT)


# ====================================================================
# SLIDE 1: Title + Project Overview
# ====================================================================
s1 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s1)
header_bar(s1, "Robotic Pick & Place for Fruit Sorting",
           "Applied Robotics (04 39984) — University of Birmingham")
footer_bar(s1, 1)

# Left column — project scope
add_rect(s1, 0.3, 1.2, 6.2, 0.4, ACCENT_BLUE)
add_textbox(s1, 0.4, 1.22, 6, 0.35, "Project Scope",
            font_size=14, bold=True)

add_rect(s1, 0.3, 1.7, 6.2, 2.4, SECTION_BG, ACCENT_BLUE)
tf = add_textbox(s1, 0.5, 1.75, 5.8, 2.3,
                 "Objective:", font_size=12, bold=True, color=ACCENT_BLUE)
add_paragraph(tf, "Sort 14 fruits into 3 baskets autonomously", 11,
              color=WHITE, space_before=4)
add_paragraph(tf, "  6 strawberries  |  3 bananas  |  5 tomatoes", 10,
              color=LIGHT_GRAY, font_name='Consolas')
add_paragraph(tf, "", 6)
add_paragraph(tf, "Operational Modes:", 12, bold=True, color=ACCENT_ORANGE)
add_paragraph(tf, "  Fully Autonomous — vision + plan + pick + place loop",
              10, color=WHITE)
add_paragraph(tf, "  Remote-Controlled — Tkinter GUI with live D415 feed",
              10, color=WHITE)
add_paragraph(tf, "", 6)
add_paragraph(tf, "Hardware:", 12, bold=True, color=ACCENT_GREEN)
add_paragraph(tf, "  Quanser QArm 4-DOF + Intel RealSense D415 RGBD", 10,
              color=WHITE)
add_paragraph(tf, "  Arm-mounted camera (pose-dependent calibration)", 10,
              color=LIGHT_GRAY)
add_paragraph(tf, "", 6)
add_paragraph(tf, "Software:", 12, bold=True, color=ACCENT_GREEN)
add_paragraph(tf, "  Python 3.13 | OpenCV | Tkinter | Vosk (voice)", 10,
              color=WHITE)

# Key features box
add_rect(s1, 0.3, 4.2, 6.2, 1.6, SECTION_BG, ACCENT_GREEN)
tf2 = add_textbox(s1, 0.5, 4.25, 5.8, 1.5,
                  "Key Features:", font_size=12, bold=True,
                  color=ACCENT_GREEN)
add_paragraph(tf2, "", 4)
add_paragraph(tf2, "HSV + shape-based fruit detection (no ML)", 10,
              color=WHITE)
add_paragraph(tf2, "Chessboard session calibration (<30s re-cal)", 10,
              color=WHITE)
add_paragraph(tf2, "13-state pick-place FSM with per-fruit tuning", 10,
              color=WHITE)
add_paragraph(tf2, "Voice control (Vosk wake-word, offline)", 10,
              color=WHITE)
add_paragraph(tf2, "60+ automated tests, quintic trajectories", 10,
              color=WHITE)

# Module info
add_rect(s1, 0.3, 5.9, 6.2, 1.1, SECTION_BG)
tf3 = add_textbox(s1, 0.5, 5.95, 5.8, 1.0,
                  "Module: Applied Robotics (04 39984)", font_size=11,
                  bold=True, color=ACCENT_BLUE)
add_paragraph(tf3, "Coordinator: Dr Amir Hajiyavand", 10, color=WHITE)
add_paragraph(tf3, "University of Birmingham — 2025/26", 10,
              color=LIGHT_GRAY)
add_paragraph(tf3, "Deadline: 1 May 2026", 10, color=ACCENT_ORANGE)

# Right column — lab photo placeholder
add_image_or_placeholder(s1, 6.8, 1.2, 6.2, 5.7,
                         "Capture.PNG",
                         "[INSERT: Photo of QArm setup in lab\n"
                         "showing arm, D415 camera, baskets, and fruits]")


# ====================================================================
# SLIDE 2: System Architecture
# ====================================================================
s2 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s2)
header_bar(s2, "System Architecture: End-to-End Pipeline",
           "From Camera Frame to Sorted Basket")
footer_bar(s2, 2)

# Pipeline boxes
pipeline = [
    ("D415\nCamera",      ACCENT_BLUE,   0.2),
    ("Calibration",       ACCENT_BLUE,   1.8),
    ("Fruit\nDetection",  ACCENT_ORANGE, 3.4),
    ("Target\nSelection", ACCENT_ORANGE, 5.0),
    ("IK\nSolver",        ACCENT_GREEN,  6.6),
    ("Quintic\nTrajectory",ACCENT_GREEN, 8.2),
    ("Pick & Place\nFSM", ACCENT_GREEN,  9.8),
    ("Basket\nSort",      ACCENT_GREEN,  11.4),
]
for label, clr, x in pipeline:
    box = add_rect(s2, x, 1.25, 1.35, 0.85, RGBColor(0x2A, 0x2A, 0x4A), clr)
    box.text_frame.word_wrap = True
    p = box.text_frame.paragraphs[0]
    p.text = label
    p.font.size = Pt(9)
    p.font.color.rgb = clr
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER

# Arrows between pipeline boxes
for i in range(len(pipeline) - 1):
    _, _, x = pipeline[i]
    add_textbox(s2, x + 1.35, 1.35, 0.45, 0.6, ">",
                font_size=18, bold=True, color=WHITE,
                alignment=PP_ALIGN.CENTER)

# 3 columns: Sensing / Planning / Actuation
col_data = [
    ("Sensing", ACCENT_BLUE, 0.2, [
        "D415 RGB + Depth (640x480)",
        "Chessboard session calibration",
        "solvePnP camera extrinsics",
        "Homography pixel-to-table",
        "Arm-mounted, pose-dependent",
        "T_cam_to_base",
    ]),
    ("Planning", ACCENT_ORANGE, 4.55, [
        "HSV colour segmentation",
        "Shape filtering (area,",
        "  circularity, aspect ratio)",
        "Ray-plane depth projection",
        "Per-fruit pick bias:",
        "  calyx dir, inscribed disk",
    ]),
    ("Actuation", ACCENT_GREEN, 8.9, [
        "4-DOF analytic IK (4 solutions)",
        "Joint-limit filter + optimal",
        "Quintic spline trajectories",
        "13-state pick-place FSM",
        "Gripper ramp + readback",
        "Per-fruit depth tuning",
    ]),
]

for title, clr, x, items in col_data:
    add_rect(s2, x, 2.35, 4.15, 0.35, clr)
    add_textbox(s2, x + 0.1, 2.37, 3.95, 0.3, title,
                font_size=13, bold=True, color=WHITE)
    add_rect(s2, x, 2.75, 4.15, 2.6, SECTION_BG, clr)
    tf = add_textbox(s2, x + 0.15, 2.8, 3.85, 2.5, items[0],
                     font_size=10, color=WHITE)
    for item in items[1:]:
        add_paragraph(tf, item, 10, color=WHITE if not item.startswith("  ")
                      else LIGHT_GRAY)

# Two-modes box
add_rect(s2, 0.2, 5.55, 6.3, 1.4, SECTION_BG, ACCENT_BLUE)
tf_modes = add_textbox(s2, 0.4, 5.6, 5.9, 1.3,
                       "Autonomous Mode:", font_size=11, bold=True,
                       color=ACCENT_BLUE)
add_paragraph(tf_modes, "Full pipeline loops: scan -> detect -> pick -> "
              "place -> repeat", 10, color=WHITE)
add_paragraph(tf_modes, "Sorts all 14 fruits without operator input", 10,
              color=LIGHT_GRAY)
add_paragraph(tf_modes, "", 6)
add_paragraph(tf_modes, "Remote Mode (GUI):", 11, bold=True,
              color=ACCENT_ORANGE)
add_paragraph(tf_modes, "Tkinter UI: live D415 feed + manual jog + basket "
              "shortcuts", 10, color=WHITE)
add_paragraph(tf_modes, "Voice control: 'robot, pick strawberry'", 10,
              color=LIGHT_GRAY)

# Software stack box
add_rect(s2, 6.7, 5.55, 6.3, 1.4, SECTION_BG, ACCENT_GREEN)
tf_sw = add_textbox(s2, 6.9, 5.6, 5.9, 1.3,
                    "Software Stack:", font_size=11, bold=True,
                    color=ACCENT_GREEN)
add_paragraph(tf_sw, "", 4)
add_paragraph(tf_sw, "Python 3.13      — all control + vision logic", 10,
              color=WHITE, font_name='Consolas')
add_paragraph(tf_sw, "OpenCV 4.x       — HSV, solvePnP, homography", 10,
              color=WHITE, font_name='Consolas')
add_paragraph(tf_sw, "Quanser SDK      — QArm driver (joint I/O)", 10,
              color=WHITE, font_name='Consolas')
add_paragraph(tf_sw, "Tkinter + Pillow — GUI with live video", 10,
              color=WHITE, font_name='Consolas')
add_paragraph(tf_sw, "Vosk + PyAudio   — offline voice commands", 10,
              color=WHITE, font_name='Consolas')


# ====================================================================
# SLIDE 3: Calibration & Vision Pipeline
# ====================================================================
s3 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s3)
header_bar(s3, "Calibration & Vision Pipeline",
           "From Pixels to Base-Frame Coordinates")
footer_bar(s3, 3)

# Left — Calibration
add_rect(s3, 0.3, 1.2, 6.2, 0.4, ACCENT_BLUE)
add_textbox(s3, 0.4, 1.22, 6, 0.35,
            "Chessboard Session Calibration", font_size=14, bold=True)

add_rect(s3, 0.3, 1.7, 6.2, 2.0, SECTION_BG, ACCENT_ORANGE)
tf_cal = add_textbox(s3, 0.5, 1.75, 5.8, 1.9,
                     "Calibration Pipeline:", font_size=11, bold=True,
                     color=ACCENT_ORANGE)
add_paragraph(tf_cal, "", 4)
add_paragraph(tf_cal, "1. Print 7x5 chessboard (30mm squares)", 10,
              color=WHITE)
add_paragraph(tf_cal, "2. Jog-touch origin corner -> chess_origin_in_base_m",
              10, color=WHITE)
add_paragraph(tf_cal, "3. Move arm to survey1 pose", 10, color=WHITE)
add_paragraph(tf_cal, "4. D415 capture -> findChessboardCorners", 10,
              color=WHITE)
add_paragraph(tf_cal, "5. solvePnP -> camera extrinsics at survey pose", 10,
              color=WHITE)
add_paragraph(tf_cal, "6. Homography -> pixel-to-table-plane projection", 10,
              color=WHITE)
add_paragraph(tf_cal, "RMS gate: warn >3mm, error >10mm", 10,
              color=ACCENT_GREEN, space_before=4)

# Key equations
add_rect(s3, 0.3, 3.8, 6.2, 1.5, SECTION_BG, ACCENT_GREEN)
tf_eq = add_textbox(s3, 0.5, 3.85, 5.8, 1.4,
                    "Key Equations:", font_size=11, bold=True,
                    color=ACCENT_GREEN)
add_paragraph(tf_eq, "", 4)
add_paragraph(tf_eq, "Rigid transform:  P_base = R . P_cam + t", 10,
              color=WHITE, font_name='Consolas')
add_paragraph(tf_eq, "Pinhole model:", 10, color=WHITE, font_name='Consolas')
add_paragraph(tf_eq, "  L.[u,v,1]' = K.[R|t].[X,Y,Z,1]'", 10,
              color=LIGHT_GRAY, font_name='Consolas')
add_paragraph(tf_eq, "Ray-plane intersection:", 10, color=WHITE,
              font_name='Consolas')
add_paragraph(tf_eq, "  L = (z_table - t_z) / (r3 . K^-1 . [u,v,1]')", 10,
              color=LIGHT_GRAY, font_name='Consolas')

# Calibration screenshot
add_image_or_placeholder(s3, 0.3, 5.4, 6.2, 1.55,
                         "figures/chessboard_7x5_30mm.png",
                         "[INSERT: Chessboard calibration image]")

# Right — Fruit Detection
add_rect(s3, 6.8, 1.2, 6.2, 0.4, ACCENT_BLUE)
add_textbox(s3, 6.9, 1.22, 6, 0.35,
            "HSV + Shape Fruit Detection", font_size=14, bold=True)

add_rect(s3, 6.8, 1.7, 6.2, 1.6, SECTION_BG, ACCENT_ORANGE)
tf_det = add_textbox(s3, 7.0, 1.75, 5.8, 1.5,
                     "Detection Pipeline:", font_size=11, bold=True,
                     color=ACCENT_ORANGE)
add_paragraph(tf_det, "", 4)
add_paragraph(tf_det, "1. Convert frame to HSV colour space", 10,
              color=WHITE)
add_paragraph(tf_det, "2. Per-fruit HSV mask (with hue-wrap for red)", 10,
              color=WHITE)
add_paragraph(tf_det, "3. Morphological open/close to clean noise", 10,
              color=WHITE)
add_paragraph(tf_det, "4. Contour extraction + area/circularity gates", 10,
              color=WHITE)
add_paragraph(tf_det, "5. Depth lookup -> ray-plane -> base-frame XYZ", 10,
              color=WHITE)

# HSV ranges table
add_rect(s3, 6.8, 3.4, 6.2, 0.3, RGBColor(0x2A, 0x2A, 0x4A))
add_textbox(s3, 6.9, 3.42, 6, 0.28,
            "Fruit       H range        S       V       Min area",
            font_size=9, bold=True, color=ACCENT_BLUE, font_name='Consolas')

hsv_rows = [
    "Banana      14-40          80+     60+     2000 px",
    "Tomato      0-10, 170-180  80+     40+     1500 px",
    "Strawberry  0-10, 170-180  80+     40+      800 px",
]
for i, row in enumerate(hsv_rows):
    bg = SECTION_BG if i % 2 == 0 else RGBColor(0x28, 0x28, 0x42)
    add_rect(s3, 6.8, 3.75 + i * 0.28, 6.2, 0.28, bg)
    add_textbox(s3, 6.9, 3.76 + i * 0.28, 6, 0.26, row,
                font_size=9, color=WHITE, font_name='Consolas')

# Strawberry special handling
add_rect(s3, 6.8, 4.65, 6.2, 1.3, SECTION_BG, ACCENT_ORANGE)
tf_straw = add_textbox(s3, 7.0, 4.7, 5.8, 1.2,
                       "Strawberry Special Handling:", font_size=11,
                       bold=True, color=ACCENT_ORANGE)
add_paragraph(tf_straw, "", 4)
add_paragraph(tf_straw, "Green calyx detection (_has_green_above)", 10,
              color=WHITE)
add_paragraph(tf_straw, "  -> distinguishes strawberry from tomato", 9,
              color=LIGHT_GRAY)
add_paragraph(tf_straw, "Widest inscribed disk replaces naive centroid", 10,
              color=WHITE)
add_paragraph(tf_straw, "  -> gripper targets fat body, not stem", 9,
              color=LIGHT_GRAY)
add_paragraph(tf_straw, "Calyx direction vector biases pick toward body", 10,
              color=WHITE)

# Detection screenshot
add_image_or_placeholder(s3, 6.8, 6.05, 6.2, 0.9,
                         "figures/live_color.png",
                         "[INSERT: Detection overlay — annotated D415 frame\n"
                         "with coloured bounding boxes per fruit type]")


# ====================================================================
# SLIDE 4: Control & Trajectory
# ====================================================================
s4 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s4)
header_bar(s4, "Control System: IK, Trajectories & FSM",
           "Inverse Kinematics | Quintic Splines | 13-State Machine")
footer_bar(s4, 4)

# Left — IK
add_rect(s4, 0.3, 1.2, 6.2, 0.4, ACCENT_BLUE)
add_textbox(s4, 0.4, 1.22, 6, 0.35,
            "Inverse Kinematics (4-DOF Analytic)", font_size=14, bold=True)

add_rect(s4, 0.3, 1.7, 6.2, 1.6, SECTION_BG, ACCENT_BLUE)
tf_ik = add_textbox(s4, 0.5, 1.75, 5.8, 1.5,
                    "4 solutions per target position:", font_size=11,
                    bold=True, color=ACCENT_BLUE)
add_paragraph(tf_ik, "", 4)
add_paragraph(tf_ik, "D1 = +sqrt(px^2+py^2): elbow-up / elbow-down", 10,
              color=WHITE, font_name='Consolas')
add_paragraph(tf_ik, "D2 = -sqrt(px^2+py^2): inverted configs", 10,
              color=WHITE, font_name='Consolas')
add_paragraph(tf_ik, "", 4)
add_paragraph(tf_ik, "Joint-limit filter: [-3pi/4, 3pi/4]", 10,
              color=ACCENT_ORANGE, font_name='Consolas')
add_paragraph(tf_ik, "Optimal = minimum joint displacement from current", 10,
              color=ACCENT_GREEN)
add_paragraph(tf_ik, "", 4)
add_paragraph(tf_ik, "DH-to-Physical mapping:", 10, bold=True,
              color=ACCENT_ORANGE)
add_paragraph(tf_ik, "  phi_1 = theta_1", 10, color=WHITE,
              font_name='Consolas')
add_paragraph(tf_ik, "  phi_2 = theta_2 + pi/2 - beta  (beta=8.13 deg)",
              10, color=WHITE, font_name='Consolas')
add_paragraph(tf_ik, "  phi_3 = theta_3 + beta", 10, color=WHITE,
              font_name='Consolas')

# Quintic Trajectory
add_rect(s4, 0.3, 3.4, 6.2, 1.5, SECTION_BG, ACCENT_GREEN)
tf_traj = add_textbox(s4, 0.5, 3.45, 5.8, 1.4,
                      "Quintic Spline Trajectory:", font_size=11,
                      bold=True, color=ACCENT_GREEN)
add_paragraph(tf_traj, "", 4)
add_paragraph(tf_traj, "s(tau) = 10*tau^3 - 15*tau^4 + 6*tau^5", 12,
              bold=True, color=WHITE, font_name='Consolas')
add_paragraph(tf_traj, "", 4)
add_paragraph(tf_traj, "Zero velocity at start and end", 10, color=WHITE)
add_paragraph(tf_traj, "Zero acceleration at start and end", 10,
              color=WHITE)
add_paragraph(tf_traj, "Jerk-continuous across segment boundaries", 10,
              color=ACCENT_GREEN)
add_paragraph(tf_traj, "", 4)
add_paragraph(tf_traj, "Improvement over cubic: eliminates audible", 10,
              color=ACCENT_ORANGE)
add_paragraph(tf_traj, "servo jerk on QArm geared drivetrain", 10,
              color=ACCENT_ORANGE)

# Trajectory plot
add_image_or_placeholder(s4, 0.3, 5.0, 6.2, 1.9,
                         "figures/trajectory_3D.png",
                         "[INSERT: 3D trajectory plot\n"
                         "showing quintic path between pick points]")

# Right — FSM
add_rect(s4, 6.8, 1.2, 6.2, 0.4, ACCENT_BLUE)
add_textbox(s4, 6.9, 1.22, 6, 0.35,
            "13-State Pick-Place FSM", font_size=14, bold=True)

add_rect(s4, 6.8, 1.7, 6.2, 2.7, SECTION_BG, ACCENT_ORANGE)
tf_fsm = add_textbox(s4, 7.0, 1.75, 5.8, 2.6,
                     "State Sequence:", font_size=11, bold=True,
                     color=ACCENT_ORANGE)
add_paragraph(tf_fsm, "", 4)
fsm_states = [
    ("INIT",           "Initialize driver + camera"),
    ("GO_HOME",        "Move to safe home position [0.45, 0, 0.49]"),
    ("SCAN",           "Capture D415 frame from survey1"),
    ("SELECT_FRUIT",   "Choose next target from detections"),
    ("APPROACH",       "Move above fruit at SAFE_Z = 0.20m"),
    ("DESCEND",        "Lower to pick height (per-fruit depth)"),
    ("CLOSE_GRIPPER",  "Ramp close + readback settle"),
    ("ASCEND_PICK",    "Lift to SAFE_Z with fruit"),
    ("MOVE_TO_BASKET", "Navigate to category basket XYZ"),
    ("DESCEND_PLACE",  "Lower into basket"),
    ("OPEN_GRIPPER",   "Release fruit"),
    ("ASCEND_PLACE",   "Lift clear, loop to SELECT_FRUIT"),
    ("DONE",           "All fruits sorted — return home"),
]
for state, desc in fsm_states:
    add_paragraph(tf_fsm, f"  {state:<17s} {desc}", 8, color=WHITE,
                  font_name='Consolas')

# Per-fruit tuning table
add_rect(s4, 6.8, 4.5, 6.2, 0.3, RGBColor(0x2A, 0x2A, 0x4A))
add_textbox(s4, 6.9, 4.52, 6, 0.28,
            "Parameter        Strawberry   Banana   Tomato",
            font_size=9, bold=True, color=ACCENT_BLUE, font_name='Consolas')

tuning_rows = [
    "Pick depth (m)   0.018        0.020    0.020",
    "X bias (cm)      -3.0         --       --",
    "Pick offset (m)  0.025        0.020    0.020",
    "Grip close       0.90         0.90     0.90",
]
for i, row in enumerate(tuning_rows):
    bg = SECTION_BG if i % 2 == 0 else RGBColor(0x28, 0x28, 0x42)
    add_rect(s4, 6.8, 4.85 + i * 0.28, 6.2, 0.28, bg)
    add_textbox(s4, 6.9, 4.86 + i * 0.28, 6, 0.26, row,
                font_size=9, color=WHITE, font_name='Consolas')

# Gripper readback box
add_rect(s4, 6.8, 6.05, 6.2, 0.9, SECTION_BG, ACCENT_GREEN)
tf_grip = add_textbox(s4, 7.0, 6.1, 5.8, 0.8,
                      "Gripper Ramp + Readback:", font_size=11, bold=True,
                      color=ACCENT_GREEN)
add_paragraph(tf_grip, "", 4)
add_paragraph(tf_grip, "Command grip_close -> wait settle -> read actual",
              10, color=WHITE)
add_paragraph(tf_grip, "Jaw adapts to fruit size (stalls at thickness)", 10,
              color=LIGHT_GRAY)
add_paragraph(tf_grip, "Confirms grasp before ascending", 10, color=WHITE)


# ====================================================================
# SLIDE 5: GUI & Voice Control
# ====================================================================
s5 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s5)
header_bar(s5, "Human Interface: GUI + Voice Control",
           "Tkinter Remote Control | Vosk Wake-Word Protocol")
footer_bar(s5, 5)

# Left — GUI
add_rect(s5, 0.3, 1.2, 6.2, 0.4, ACCENT_BLUE)
add_textbox(s5, 0.4, 1.22, 6, 0.35,
            "Tkinter GUI — Remote Control", font_size=14, bold=True)

# Layout diagram
add_rect(s5, 0.3, 1.7, 6.2, 2.3, SECTION_BG, ACCENT_BLUE)
tf_gui = add_textbox(s5, 0.5, 1.75, 5.8, 2.2,
                     "GUI Layout:", font_size=11, bold=True,
                     color=ACCENT_BLUE)
add_paragraph(tf_gui, "", 4)
add_paragraph(tf_gui,
              "+------------------+-------------------------+", 8,
              color=ACCENT_BLUE, font_name='Consolas')
add_paragraph(tf_gui,
              "| Live D415 feed   | EMERGENCY STOP          |", 8,
              color=WHITE, font_name='Consolas')
add_paragraph(tf_gui,
              "| (640 x 360)      | Manual Teleop (XYZ +/-) |", 8,
              color=WHITE, font_name='Consolas')
add_paragraph(tf_gui,
              "|  + detection     | Basket: B / T / S       |", 8,
              color=WHITE, font_name='Consolas')
add_paragraph(tf_gui,
              "|    overlay       | Gripper: open / close    |", 8,
              color=WHITE, font_name='Consolas')
add_paragraph(tf_gui,
              "|                  | Mode: Manual | Auto      |", 8,
              color=WHITE, font_name='Consolas')
add_paragraph(tf_gui,
              "|                  | Status: ...              |", 8,
              color=WHITE, font_name='Consolas')
add_paragraph(tf_gui,
              "+------------------+-------------------------+", 8,
              color=ACCENT_BLUE, font_name='Consolas')

# GUI features
add_rect(s5, 0.3, 4.1, 6.2, 1.5, SECTION_BG, ACCENT_GREEN)
tf_feat = add_textbox(s5, 0.5, 4.15, 5.8, 1.4,
                      "Key Features:", font_size=11, bold=True,
                      color=ACCENT_GREEN)
add_paragraph(tf_feat, "", 4)
add_paragraph(tf_feat,
              "Live camera feed updates during picks (not frozen)", 10,
              color=WHITE)
add_paragraph(tf_feat,
              "E-STOP halts arm + logs EE pose (not just a flag)", 10,
              color=WHITE)
add_paragraph(tf_feat,
              "Concurrent-read safe (D415 driver not thread-safe)", 10,
              color=WHITE)
add_paragraph(tf_feat,
              "Detection overlay: coloured boxes per fruit type", 10,
              color=WHITE)
add_paragraph(tf_feat,
              "Auto batch: press 'all' to sort every fruit", 10,
              color=WHITE)

# GUI screenshot placeholder
add_image_or_placeholder(s5, 0.3, 5.7, 6.2, 1.25,
                         "1Capture.PNG",
                         "[INSERT: GUI screenshot showing live feed\n"
                         "with detection overlay and control panel]")

# Right — Voice Control
add_rect(s5, 6.8, 1.2, 6.2, 0.4, ACCENT_BLUE)
add_textbox(s5, 6.9, 1.22, 6, 0.35,
            "Voice Control — Vosk Wake-Word", font_size=14, bold=True)

# Wake-word protocol
add_rect(s5, 6.8, 1.7, 6.2, 2.0, SECTION_BG, ACCENT_ORANGE)
tf_voice = add_textbox(s5, 7.0, 1.75, 5.8, 1.9,
                       "Wake-Word Protocol:", font_size=11, bold=True,
                       color=ACCENT_ORANGE)
add_paragraph(tf_voice, "", 4)
add_paragraph(tf_voice,
              '1. Hear "robot" -> enter awaiting phase (2s timer)', 10,
              color=WHITE)
add_paragraph(tf_voice,
              "2. Within 2s: hear valid command -> dispatch", 10,
              color=WHITE)
add_paragraph(tf_voice,
              "3. Timeout or noise -> return to idle", 10, color=WHITE)
add_paragraph(tf_voice,
              '"stop" always dispatched (no wake needed)', 10,
              color=ACCENT_GREEN)
add_paragraph(tf_voice, "", 4)
add_paragraph(tf_voice, "Example: 'robot, pick strawberry'", 11, bold=True,
              color=WHITE)
add_paragraph(tf_voice, "Example: 'stop' (immediate halt)", 11, bold=True,
              color=ACCENT_RED)

# Grammar + Architecture
add_rect(s5, 6.8, 3.8, 6.2, 1.4, SECTION_BG, ACCENT_GREEN)
tf_arch = add_textbox(s5, 7.0, 3.85, 5.8, 1.3,
                      "Architecture:", font_size=11, bold=True,
                      color=ACCENT_GREEN)
add_paragraph(tf_arch, "", 4)
add_paragraph(tf_arch,
              "Vosk + PyAudio — fully offline (no cloud)", 10, color=WHITE)
add_paragraph(tf_arch,
              "Small English model (vosk-model-small-en-us-0.15)", 9,
              color=LIGHT_GRAY)
add_paragraph(tf_arch,
              "Threaded: VoiceController daemon + callback dispatch", 10,
              color=WHITE)
add_paragraph(tf_arch,
              "Graceful degradation: GUI works without mic/model", 10,
              color=ACCENT_ORANGE)
add_paragraph(tf_arch,
              "Sample rate: 16 kHz, chunk: 4000 frames", 10,
              color=LIGHT_GRAY)

# Grammar table
add_rect(s5, 6.8, 5.3, 6.2, 0.3, RGBColor(0x2A, 0x2A, 0x4A))
add_textbox(s5, 6.9, 5.32, 6, 0.28,
            "Grammar Words (Vosk restricted vocabulary):",
            font_size=9, bold=True, color=ACCENT_BLUE)

add_rect(s5, 6.8, 5.65, 6.2, 0.55, SECTION_BG)
add_textbox(s5, 6.9, 5.67, 6, 0.5,
            "robot | strawberry | banana | tomato | all | stop | "
            "home | refresh | survey",
            font_size=10, color=WHITE, font_name='Consolas')

# Theory link
add_rect(s5, 6.8, 6.3, 6.2, 0.7, SECTION_BG, ACCENT_BLUE)
tf_theory = add_textbox(s5, 7.0, 6.35, 5.8, 0.6,
                        "Theory Link:", font_size=11, bold=True,
                        color=ACCENT_BLUE)
add_paragraph(tf_theory,
              "FSM pattern: same design as pick-place controller", 10,
              color=WHITE)
add_paragraph(tf_theory,
              "Real-time audio processing at 16kHz (signal processing)", 10,
              color=WHITE)


# ====================================================================
# SLIDE 6: Results & Validation
# ====================================================================
s6 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s6)
header_bar(s6, "Experimental Results & Validation",
           "Calibration | Detection | Pick & Place | Test Suite")
footer_bar(s6, 6)

# Top row: 3 result boxes
# Calibration
add_rect(s6, 0.2, 1.2, 4.2, 2.2, SECTION_BG, ACCENT_BLUE)
tf_r1 = add_textbox(s6, 0.4, 1.25, 3.8, 2.1,
                    "Calibration", font_size=13, bold=True,
                    color=ACCENT_BLUE)
add_paragraph(tf_r1, "", 4)
add_paragraph(tf_r1, "Chessboard RMS: ~36 mm", 10, color=WHITE)
add_paragraph(tf_r1, "Session-cal refresh: <30 s", 10, color=WHITE)
add_paragraph(tf_r1, "solvePnP extrinsics verified", 10, color=WHITE)
add_paragraph(tf_r1, "survey1 pose repeatable", 10, color=WHITE)
add_paragraph(tf_r1, "Settle tolerance: 0.04 rad", 10, color=LIGHT_GRAY)
add_paragraph(tf_r1, "Downstream residual gate: 10mm", 10,
              color=LIGHT_GRAY)

# Detection
add_rect(s6, 4.6, 1.2, 4.2, 2.2, SECTION_BG, ACCENT_ORANGE)
tf_r2 = add_textbox(s6, 4.8, 1.25, 3.8, 2.1,
                    "Fruit Detection", font_size=13, bold=True,
                    color=ACCENT_ORANGE)
add_paragraph(tf_r2, "", 4)
add_paragraph(tf_r2, "3 fruit types classified", 10, color=WHITE)
add_paragraph(tf_r2, "HSV + shape pipeline (no ML)", 10, color=WHITE)
add_paragraph(tf_r2, "Depth projection to base frame", 10, color=WHITE)
add_paragraph(tf_r2, "Calyx-based strawberry ID", 10, color=WHITE)
add_paragraph(tf_r2, "Hue-wrap handles red fruits", 10, color=LIGHT_GRAY)
add_paragraph(tf_r2, "Area/circularity gates per type", 10,
              color=LIGHT_GRAY)

# Pick & Place
add_rect(s6, 9.0, 1.2, 4.2, 2.2, SECTION_BG, ACCENT_GREEN)
tf_r3 = add_textbox(s6, 9.2, 1.25, 3.8, 2.1,
                    "Pick & Place", font_size=13, bold=True,
                    color=ACCENT_GREEN)
add_paragraph(tf_r3, "", 4)
add_paragraph(tf_r3, "14 fruits -> 3 baskets", 10, color=WHITE)
add_paragraph(tf_r3, "Autonomous full-cycle demo", 10, color=WHITE)
add_paragraph(tf_r3, "Per-fruit depth tuning", 10, color=WHITE)
add_paragraph(tf_r3, "Gripper readback confirms grasp", 10, color=WHITE)
add_paragraph(tf_r3, "IK fallback for edge poses", 10, color=LIGHT_GRAY)
add_paragraph(tf_r3, "Quintic smooth motion", 10, color=LIGHT_GRAY)

# Test suite summary
add_rect(s6, 0.2, 3.6, 12.8, 0.35, ACCENT_GREEN)
add_textbox(s6, 0.4, 3.62, 12.4, 0.3,
            "Test Suite: 60/60 passing across 6 suites + voice control",
            font_size=13, bold=True, color=WHITE)

test_suites = [
    "test_integration         — end-to-end FSM mock tests",
    "test_fruit_detector      — HSV + shape pipeline unit tests",
    "test_calibrate_chessboard — solvePnP + homography tests",
    "test_pick_single         — single-fruit pick cycle tests",
    "test_survey_capture      — camera capture + settle tests",
    "test_picker_viewer       — UI event loop + live feed tests",
    "test_voice_control       — parser + dispatch integration",
]
add_rect(s6, 0.2, 4.0, 12.8, 1.6, SECTION_BG)
tf_tests = add_textbox(s6, 0.4, 4.05, 12.4, 1.5, test_suites[0],
                       font_size=9, color=WHITE, font_name='Consolas')
for suite in test_suites[1:]:
    add_paragraph(tf_tests, suite, 9, color=WHITE, font_name='Consolas')

# Bottom: plots from figures/
add_image_or_placeholder(s6, 0.2, 5.75, 4.2, 1.2,
                         "figures/joint_angles_time.png",
                         "[INSERT: Joint angles over pick cycle\n"
                         "from figures/joint_angles_time.png]")

add_image_or_placeholder(s6, 4.6, 5.75, 4.2, 1.2,
                         "figures/gripper_state.png",
                         "[INSERT: Gripper commanded vs actual\n"
                         "from figures/gripper_state.png]")

add_image_or_placeholder(s6, 9.0, 5.75, 4.2, 1.2,
                         "figures/top_view.png",
                         "[INSERT: Top-view trajectory\n"
                         "from figures/top_view.png]")


# ====================================================================
# SLIDE 7: Challenges, Solutions & Feedback
# ====================================================================
s7 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s7)
header_bar(s7, "Challenges, Solutions & Improvements",
           "Problems Encountered | How We Solved Them | Feedback Response")
footer_bar(s7, 7)

# Left — 6 challenge/solution pairs
add_rect(s7, 0.3, 1.2, 6.2, 0.4, ACCENT_ORANGE)
add_textbox(s7, 0.4, 1.22, 6, 0.35,
            "Challenges & Solutions", font_size=14, bold=True)

challenges = [
    ("Calibration drifts when QArm base is bumped",
     "Session-cal system: quick re-cal in <30s with chessboard residual "
     "gates (3mm warn, 10mm error)"),
    ("Strawberry grip too loose at 0.72",
     "Reverted to 0.90 + gripper ramp+readback adapts jaw to fruit size "
     "(stalls at thickness)"),
    ("Tomato detection false negatives",
     "Fixed cam_extrinsics_survey1 passthrough bug + HSV hue-wrap for red "
     "range (0-10, 170-180)"),
    ("Concurrent D415 reads crash driver",
     "Camera reader auto-pauses during capture_fruits(); Quanser Video3D "
     "is not thread-safe"),
    ("Settle tolerance too tight (0.015 rad)",
     "Relaxed to 0.04 rad; downstream chessboard residual (10mm) is the "
     "real mm-level guard"),
    ("Strawberry centroid off-centre (calyx)",
     "Widest inscribed disk centroid + calyx direction vector biases pick "
     "toward fat body"),
]

y = 1.7
for i, (chal, sol) in enumerate(challenges):
    add_rect(s7, 0.3, y, 6.2, 0.8, SECTION_BG,
             ACCENT_ORANGE if i % 2 == 0 else ACCENT_GREEN)
    tf_c = add_textbox(s7, 0.5, y + 0.02, 5.8, 0.75,
                       f"Challenge: {chal}", font_size=9, bold=True,
                       color=ACCENT_ORANGE)
    add_paragraph(tf_c, f"Solution: {sol}", 9, color=ACCENT_GREEN)
    y += 0.85

# Right — Improvements from feedback
add_rect(s7, 6.8, 1.2, 6.2, 0.4, ACCENT_GREEN)
add_textbox(s7, 6.9, 1.22, 6, 0.35,
            "Improvements from Feedback", font_size=14, bold=True)

add_rect(s7, 6.8, 1.7, 6.2, 4.0, SECTION_BG, ACCENT_GREEN)
tf_imp = add_textbox(s7, 7.0, 1.75, 5.8, 3.9,
                     "From A1/A2 Feedback:", font_size=12, bold=True,
                     color=ACCENT_GREEN)

improvements = [
    ("Quantitative analysis",
     "Added RMS residuals, settling times, success rates — "
     "not just qualitative descriptions"),
    ("Quintic replaces cubic",
     "Zero acceleration at boundaries eliminates servo jerk; "
     "smoother, quieter motion on geared drivetrain"),
    ("E-STOP with actual arm halt",
     "E-STOP now freezes joints + logs EE pose, "
     "not just setting a flag the FSM checks later"),
    ("Voice control added",
     "Hands-free operation via Vosk wake-word; "
     "accessibility + safety (eyes on arm, not keyboard)"),
    ("60+ automated regression tests",
     "6 test suites + voice tests; prevents regressions "
     "across 178 commits of iterative tuning"),
    ("Live camera feed during picks",
     "D415 frame updates in real-time during pick motion; "
     "operator sees what the arm sees, not a stale snapshot"),
]

for title, detail in improvements:
    add_paragraph(tf_imp, "", 4)
    add_paragraph(tf_imp, f"  {title}", 10, bold=True, color=WHITE)
    add_paragraph(tf_imp, f"  {detail}", 9, color=LIGHT_GRAY)

# Bottom-right: development stats
add_rect(s7, 6.8, 5.85, 6.2, 1.1, SECTION_BG, ACCENT_BLUE)
tf_stats = add_textbox(s7, 7.0, 5.9, 5.8, 1.0,
                       "Development Stats:", font_size=11, bold=True,
                       color=ACCENT_BLUE)
add_paragraph(tf_stats, "", 4)
add_paragraph(tf_stats, "178 commits across 7 branches", 10, color=WHITE,
              font_name='Consolas')
add_paragraph(tf_stats, "60+ automated tests (all passing)", 10,
              color=WHITE, font_name='Consolas')
add_paragraph(tf_stats, "~20 Python modules (driver, vision, control, UI)",
              10, color=WHITE, font_name='Consolas')
add_paragraph(tf_stats, "Full CI: preflight -> calibrate -> sort -> report",
              10, color=LIGHT_GRAY, font_name='Consolas')


# ====================================================================
# SLIDE 8: Conclusion & Future Work
# ====================================================================
s8 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s8)
header_bar(s8, "Conclusion & Future Work",
           "Summary | Theory-Practice Bridge | Next Steps")
footer_bar(s8, 8)

# Left — Conclusion
add_rect(s8, 0.3, 1.2, 4.0, 0.4, ACCENT_BLUE)
add_textbox(s8, 0.4, 1.22, 3.8, 0.35,
            "Conclusion", font_size=14, bold=True)

add_rect(s8, 0.3, 1.7, 4.0, 3.6, SECTION_BG, ACCENT_GREEN)
tf_conc = add_textbox(s8, 0.5, 1.75, 3.6, 3.5,
                      "Achievements:", font_size=12, bold=True,
                      color=ACCENT_GREEN)
add_paragraph(tf_conc, "", 4)
add_paragraph(tf_conc,
              "Complete autonomous fruit sorting system", 10, color=WHITE)
add_paragraph(tf_conc,
              "14 fruits (3 types) sorted into 3 baskets", 10, color=WHITE)
add_paragraph(tf_conc, "", 4)
add_paragraph(tf_conc,
              "Two operational modes:", 10, bold=True, color=ACCENT_BLUE)
add_paragraph(tf_conc,
              "  Fully autonomous (vision loop)", 10, color=WHITE)
add_paragraph(tf_conc,
              "  Remote GUI with live feed", 10, color=WHITE)
add_paragraph(tf_conc, "", 4)
add_paragraph(tf_conc,
              "Voice control for hands-free", 10, bold=True,
              color=ACCENT_ORANGE)
add_paragraph(tf_conc,
              "  operation and accessibility", 10, color=LIGHT_GRAY)
add_paragraph(tf_conc, "", 4)
add_paragraph(tf_conc,
              "Theory-practice bridge:", 10, bold=True, color=ACCENT_GREEN)
add_paragraph(tf_conc,
              "  IK, trajectory gen, visual", 10, color=WHITE)
add_paragraph(tf_conc,
              "  servoing, FSM — all applied", 10, color=WHITE)
add_paragraph(tf_conc,
              "  from lecture to real hardware", 10, color=WHITE)

# Centre — Theory-Practice table
add_rect(s8, 4.5, 1.2, 5.0, 0.4, ACCENT_BLUE)
add_textbox(s8, 4.6, 1.22, 4.8, 0.35,
            "Theory-Practice Bridge", font_size=14, bold=True)

add_rect(s8, 4.5, 1.7, 5.0, 0.35, RGBColor(0x2A, 0x2A, 0x4A))
add_textbox(s8, 4.6, 1.72, 4.8, 0.3,
            "Concept            Our Implementation",
            font_size=10, bold=True, color=ACCENT_BLUE, font_name='Consolas')

theory_rows = [
    ("Inverse Kinematics",
     "4-solution analytic solver",
     "Spong & Vidyasagar"),
    ("Trajectory Gen",
     "Quintic splines (zero jerk)",
     "Craig Ch.7"),
    ("Visual Servoing",
     "D415 solvePnP + homography",
     "Corke: eye-in-hand"),
    ("State Machines",
     "13-state pick-place FSM",
     "FSM / Stateflow theory"),
]

for i, (concept, impl, ref) in enumerate(theory_rows):
    bg = SECTION_BG if i % 2 == 0 else RGBColor(0x28, 0x28, 0x42)
    add_rect(s8, 4.5, 2.1 + i * 0.85, 5.0, 0.8, bg, ACCENT_BLUE)
    tf_row = add_textbox(s8, 4.65, 2.12 + i * 0.85, 4.7, 0.75,
                         concept, font_size=11, bold=True,
                         color=ACCENT_ORANGE)
    add_paragraph(tf_row, impl, 10, color=WHITE)
    add_paragraph(tf_row, f"Ref: {ref}", 9, color=LIGHT_GRAY)

# Right — Future Work
add_rect(s8, 9.7, 1.2, 3.3, 0.4, ACCENT_RED)
add_textbox(s8, 9.8, 1.22, 3.1, 0.35,
            "Future Work", font_size=14, bold=True)

add_rect(s8, 9.7, 1.7, 3.3, 3.6, SECTION_BG, ACCENT_RED)
tf_fw = add_textbox(s8, 9.85, 1.75, 3.0, 3.5,
                    "Recommendations:", font_size=11, bold=True,
                    color=ACCENT_RED)

future_items = [
    ("Full Parallax Projection",
     "Eliminate pose-dependent calibration limitation"),
    ("ML-Based Detection",
     "Replace hand-tuned HSV with trained classifier"),
    ("Soft Gripper Design",
     "Custom compliant fingers for delicate fruits"),
    ("Multi-Arm Coordination",
     "Parallel sorting for higher throughput"),
]

for title, desc in future_items:
    add_paragraph(tf_fw, "", 4)
    add_paragraph(tf_fw, title, 10, bold=True, color=WHITE)
    add_paragraph(tf_fw, desc, 9, color=LIGHT_GRAY)

# Bottom — Video placeholder
add_image_or_placeholder(s8, 0.3, 5.5, 12.7, 1.5,
                         "figures/ee_position_time.png",
                         "[INSERT: Video/GIF of full autonomous sorting "
                         "cycle\nshowing scan -> detect -> pick -> place "
                         "for all 14 fruits]")


# ====================================================================
# SAVE
# ====================================================================
_out = os.path.join(os.path.dirname(os.path.abspath(__file__)),
                    "FinalPresentation_FruitSorting.pptx")
prs.save(_out)
print(f"Saved to: {_out}")
